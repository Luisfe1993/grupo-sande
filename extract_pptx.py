from pptx import Presentation
import os

pptx_path = r'c:\Users\luissande\projects\grupo-sande\2507 SUCESION 6 (1).pptx'
img_dir = r'c:\Users\luissande\projects\grupo-sande\pptx_extract'
os.makedirs(img_dir, exist_ok=True)

prs = Presentation(pptx_path)
img_count = 0

for slide_idx, slide in enumerate(prs.slides):
    slide_num = slide_idx + 1
    print(f"\n=== SLIDE {slide_num} ===")
    
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    print(f"  TEXT: {text}")
        
        if shape.shape_type == 13:  # Picture
            img_count += 1
            try:
                image = shape.image
                ext = image.content_type.split('/')[-1].replace('jpeg', 'jpg')
                fname = f"slide{slide_num}_img{img_count}.{ext}"
                fpath = os.path.join(img_dir, fname)
                with open(fpath, 'wb') as f:
                    f.write(image.blob)
                print(f"  IMAGE: {fname}")
            except ValueError:
                print(f"  IMAGE: (linked, not embedded - skipped)")

print(f"\nTotal slides: {len(prs.slides)}")
print(f"Total images extracted: {img_count}")
